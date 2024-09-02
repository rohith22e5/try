import os
import shutil
import sys
import re
from collections import defaultdict
from PyQt5.QtWidgets import (QApplication, QMainWindow, QFileDialog, QVBoxLayout, QHBoxLayout, QPushButton, 
                             QWidget, QListWidget, QListWidgetItem, QProgressBar, QMessageBox, QLabel)
from PyQt5.QtGui import QIcon, QFont
from PyQt5.QtCore import Qt
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from sklearn.metrics import silhouette_score
import logging
import nltk
from bs4 import BeautifulSoup

# Set up logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

# Download necessary NLTK data
nltk.download('stopwords')
nltk.download('punkt')
stopwords = nltk.corpus.stopwords.words('english')

class DocumentOrganizer(QMainWindow):
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        self.setWindowTitle('Automated Document Organizer')
        self.setGeometry(100, 100, 800, 600)
        self.setWindowIcon(QIcon('icon.png'))  # Set an icon for the window

        self.layout = QVBoxLayout()
        self.button_layout = QHBoxLayout()

        self.title_label = QLabel('Automated Document Organizer', self)
        self.title_label.setAlignment(Qt.AlignCenter)
        self.title_label.setFont(QFont('Arial', 20))
        self.layout.addWidget(self.title_label)

        self.btn_select_folder = QPushButton('Select Input Folder', self)
        self.btn_select_folder.clicked.connect(self.open_folder_dialog)
        self.btn_select_folder.setFont(QFont('Arial', 12))
        self.button_layout.addWidget(self.btn_select_folder)
        
        self.btn_select_output_folder = QPushButton('Select Output Folder', self)
        self.btn_select_output_folder.clicked.connect(self.open_output_folder_dialog)
        self.btn_select_output_folder.setFont(QFont('Arial', 12))
        self.button_layout.addWidget(self.btn_select_output_folder)

        self.btn_add_file = QPushButton('Add Individual File', self)
        self.btn_add_file.clicked.connect(self.add_file_dialog)
        self.btn_add_file.setFont(QFont('Arial', 12))
        self.button_layout.addWidget(self.btn_add_file)
        
        self.layout.addLayout(self.button_layout)

        self.input_folder_label = QLabel('', self)
        self.input_folder_label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.input_folder_label)
        
        self.result_list = QListWidget(self)
        self.layout.addWidget(self.result_list)
        
        self.progress_bar = QProgressBar(self)
        self.layout.addWidget(self.progress_bar)
        
        self.status_label = QLabel('', self)
        self.status_label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.status_label)

        container = QWidget()
        container.setLayout(self.layout)
        self.setCentralWidget(container)
        
        self.input_folder = ""
        self.output_folder = ""

        self.setStyleSheet("""
            QMainWindow {
                background-color: #f0f0f0;
            }
            QPushButton {
                background-color: #4CAF50; 
                color: white;
                border: none;
                padding: 10px 24px;
                text-align: center;
                text-decoration: none;
                display: inline-block;
                font-size: 16px;
                margin: 4px 2px;
                cursor: pointer;
                border-radius: 12px;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
            QLabel {
                color: #333;
            }
            QProgressBar {
                border: 2px solid #4CAF50;
                border-radius: 5px;
                text-align: center;
            }
            QProgressBar::chunk {
                background-color: #4CAF50;
                width: 20px;
            }
            QListWidget {
                background-color: #ffffff;
                border: 1px solid #ddd;
                border-radius: 5px;
            }
        """)

    def open_folder_dialog(self):
        self.input_folder = QFileDialog.getExistingDirectory(self, "Select Input Folder")
        if self.input_folder:
            self.input_folder_label.setText(f"Input Folder: {self.input_folder}")
    
    def open_output_folder_dialog(self):
        self.output_folder = QFileDialog.getExistingDirectory(self, "Select Output Folder")
        if self.output_folder:
            self.status_label.setText(f"Output Folder: {self.output_folder}")
            if self.input_folder:
                self.organize_input_folder()

    def add_file_dialog(self):
        options = QFileDialog.Options()
        options |= QFileDialog.DontUseNativeDialog
        file_path, _ = QFileDialog.getOpenFileName(self, "Select File to Add", "", "All Files (*);;PDF Files (*.pdf);;Word Files (*.docx);;PowerPoint Files (*.pptx);;Text Files (*.txt)", options=options)
        if file_path:
            if not self.input_folder:
                QMessageBox.warning(self, "Input Folder Not Selected", "Please select an input folder first.")
                return
            if not self.output_folder:
                QMessageBox.warning(self, "Output Folder Not Selected", "Please select an output folder first.")
                return
            self.process_individual_file(file_path)

    def get_files_from_folder(self, folder):
        supported_extensions = ('.pdf', '.docx', '.pptx', '.txt')
        files = []
        for root, _, filenames in os.walk(folder):
            for filename in filenames:
                if filename.lower().endswith(supported_extensions):
                    files.append(os.path.join(root, filename))
        logging.debug(f"Files found: {files}")
        return files

    def is_folder_categorized(self, folder):
        subfolders = [f.path for f in os.scandir(folder) if f.is_dir()]
        if not subfolders:
            return False
        
        for subfolder in subfolders:
            files = self.get_files_from_folder(subfolder)
            if files:
                return True
        return False
    
    def organize_input_folder(self):
        if self.is_folder_categorized(self.input_folder):
            shutil.copytree(self.input_folder, self.output_folder, dirs_exist_ok=True)
            self.status_label.setText("Folders are already categorized and copied.")
        else:
            files = self.get_files_from_folder(self.input_folder)
            if files:
                self.process_files(files, self.output_folder)
            else:
                QMessageBox.warning(self, "No Files", "No supported files found in the selected directory.")

    def process_files(self, files, output_dir):
        self.progress_bar.setValue(0)
        self.progress_bar.setMaximum(len(files))
        
        categorized_files = self.categorize_files(files, self.progress_bar)
        self.organize_files_into_folders(categorized_files, output_dir, self.progress_bar)
        self.display_results(categorized_files)
        
    def process_individual_file(self, file_path):
        if self.is_folder_categorized(self.input_folder):
            folder_to_copy = self.categorize_individual_file(file_path)
            shutil.copy(file_path, os.path.join(self.output_folder, folder_to_copy))
            self.status_label.setText(f"File added to {folder_to_copy} in output directory.")
        else:
            self.status_label.setText("Input folder is not categorized. Please categorize the input folder first.")

    def categorize_individual_file(self, file_path):
        categorized_folders = [f.name for f in os.scandir(self.input_folder) if f.is_dir()]
        text = self.extract_text_from_file(file_path)
        if not text or len(text.split()) < 3:
            return "Uncategorized"
        
        vectorizer = TfidfVectorizer(stop_words='english')
        try:
            X = vectorizer.fit_transform([text])
        except ValueError:
            return "Uncategorized"
        
        best_match_folder = "Uncategorized"
        best_silhouette_score = -1
        
        for folder in categorized_folders:
            folder_files = self.get_files_from_folder(os.path.join(self.input_folder, folder))
            folder_texts = [self.extract_text_from_file(file) for file in folder_files]
            folder_texts = [t for t in folder_texts if t and len(t.split()) >= 3]
            if not folder_texts:
                continue
            
            try:
                folder_X = vectorizer.fit_transform(folder_texts)
            except ValueError:
                continue
            
            combined_X = vectorizer.fit_transform(folder_texts + [text])
            labels = KMeans(n_clusters=2, random_state=42).fit_predict(combined_X)
            silhouette_avg = silhouette_score(combined_X, labels)
            
            if silhouette_avg > best_silhouette_score:
                best_silhouette_score = silhouette_avg
                best_match_folder = folder
        
        if best_match_folder == "Uncategorized":
            best_match_folder = "Category_" + str(len(categorized_folders) + 1)
            os.makedirs(os.path.join(self.output_folder, best_match_folder), exist_ok=True)
        
        return best_match_folder
        
    def display_results(self, categorized_files):
        self.result_list.clear()
        for category, files in categorized_files.items():
            category_item = QListWidgetItem(f"Category: {category}")
            category_item.setFont(QFont('Arial', 14, QFont.Bold))
            self.result_list.addItem(category_item)
            for file in files:
                self.result_list.addItem(f" - {os.path.basename(file)}")
        self.status_label.setText("Document organization completed.")

    def extract_text_from_file(self, file_path):
        logging.debug(f"Extracting text from: {file_path}")
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        if ext == '.pdf':
            return self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return self.extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return self.extract_text_from_pptx(file_path)
        elif ext == '.txt':
            return self.extract_text_from_txt(file_path)
        return ""

    def extract_text_from_pdf(self, file_path):
        try:
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        except Exception as e:
            logging.error(f"Error extracting text from PDF: {e}")
            return ""

    def extract_text_from_docx(self, file_path):
        try:
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logging.error(f"Error extracting text from DOCX: {e}")
            return ""

    def extract_text_from_pptx(self, file_path):
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            return text
        except Exception as e:
            logging.error(f"Error extracting text from PPTX: {e}")
            return ""

    def extract_text_from_txt(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logging.error(f"Error extracting text from TXT: {e}")
            return ""

    def categorize_files(self, files, progress_bar):
        texts = [self.extract_text_from_file(file) for file in files]
        vectorizer = TfidfVectorizer(stop_words='english')
        X = vectorizer.fit_transform(texts)

        kmeans = KMeans(n_clusters=5, random_state=42)
        labels = kmeans.fit_predict(X)

        categorized_files = defaultdict(list)
        for label, file in zip(labels, files):
            categorized_files[f"Category_{label + 1}"].append(file)
        
        progress_bar.setValue(len(files))
        return categorized_files

    def organize_files_into_folders(self, categorized_files, output_dir, progress_bar):
        os.makedirs(output_dir, exist_ok=True)
        total_files = sum(len(files) for files in categorized_files.values())
        processed_files = 0
        
        for category, files in categorized_files.items():
            category_dir = os.path.join(output_dir, category)
            os.makedirs(category_dir, exist_ok=True)
            for file in files:
                shutil.copy(file, category_dir)
                processed_files += 1
                progress_bar.setValue(processed_files)

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = DocumentOrganizer()
    ex.show()
    sys.exit(app.exec_())
