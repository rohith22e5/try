import os
import shutil
import sys
from collections import defaultdict
from PyQt5.QtWidgets import (QApplication, QMainWindow, QFileDialog, QVBoxLayout, QHBoxLayout, QPushButton, 
                             QWidget, QListWidget, QListWidgetItem, QProgressBar, QMessageBox, QLabel)
from PyQt5.QtGui import QIcon, QFont
from PyQt5.QtCore import Qt
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from sklearn.metrics import silhouette_score
import logging
import nltk
from easyocr import Reader
from ultralytics import YOLO

# Set up logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

# Download necessary NLTK data
nltk.download('stopwords')
nltk.download('punkt')
stopwords = nltk.corpus.stopwords.words('english')

class DocumentOrganizer(QMainWindow):
    def __init__(self):
        super().__init__()
        self.initUI()
        self.reader = Reader(['en'])  # Initialize EasyOCR reader
        self.yolo_model = YOLO('best.pt')  # Load YOLOv8 model

    def initUI(self):
        self.setWindowTitle('Automated Document Organizer')
        self.setGeometry(100, 100, 800, 600)
        self.setWindowIcon(QIcon('icon.png'))  # Set an icon for the window

        self.layout = QVBoxLayout()
        self.button_layout = QHBoxLayout()

        self.title_label = QLabel('Automated Document Organizer', self)
        self.title_label.setAlignment(Qt.AlignCenter)
        self.title_label.setFont(QFont('Arial', 20))
        self.layout.addWidget(self.title_label)

        self.btn_select_folder = QPushButton('Select Input Folder', self)
        self.btn_select_folder.clicked.connect(self.open_folder_dialog)
        self.btn_select_folder.setFont(QFont('Arial', 12))
        self.button_layout.addWidget(self.btn_select_folder)
        
        self.btn_select_output_folder = QPushButton('Select Output Folder', self)
        self.btn_select_output_folder.clicked.connect(self.open_output_folder_dialog)
        self.btn_select_output_folder.setFont(QFont('Arial', 12))
        self.button_layout.addWidget(self.btn_select_output_folder)

        self.btn_add_file = QPushButton('Add Individual File', self)
        self.btn_add_file.clicked.connect(self.add_file_dialog)
        self.btn_add_file.setFont(QFont('Arial', 12))
        self.button_layout.addWidget(self.btn_add_file)
        
        self.layout.addLayout(self.button_layout)

        self.input_folder_label = QLabel('', self)
        self.input_folder_label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.input_folder_label)
        
        self.result_list = QListWidget(self)
        self.layout.addWidget(self.result_list)
        
        self.progress_bar = QProgressBar(self)
        self.layout.addWidget(self.progress_bar)
        
        self.status_label = QLabel('', self)
        self.status_label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.status_label)

        container = QWidget()
        container.setLayout(self.layout)
        self.setCentralWidget(container)
        
        self.input_folder = ""
        self.output_folder = ""

        self.setStyleSheet("""
            QMainWindow {
                background-color: #f0f0f0;
            }
            QPushButton {
                background-color: #4CAF50; 
                color: white;
                border: none;
                padding: 10px 24px;
                text-align: center;
                text-decoration: none;
                display: inline-block;
                font-size: 16px;
                margin: 4px 2px;
                cursor: pointer;
                border-radius: 12px;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
            QLabel {
                color: #333;
            }
            QProgressBar {
                border: 2px solid #4CAF50;
                border-radius: 5px;
                text-align: center;
            }
            QProgressBar::chunk {
                background-color: #4CAF50;
                width: 20px;
            }
            QListWidget {
                background-color: #ffffff;
                border: 1px solid #ddd;
                border-radius: 5px;
            }
        """)

    def open_folder_dialog(self):
        self.input_folder = QFileDialog.getExistingDirectory(self, "Select Input Folder")
        if self.input_folder:
            self.input_folder_label.setText(f"Input Folder: {self.input_folder}")
    
    def open_output_folder_dialog(self):
        self.output_folder = QFileDialog.getExistingDirectory(self, "Select Output Folder")
        if self.output_folder:
            self.status_label.setText(f"Output Folder: {self.output_folder}")
            if self.input_folder:
                self.organize_input_folder()

    def add_file_dialog(self):
        options = QFileDialog.Options()
        options |= QFileDialog.DontUseNativeDialog
        file_path, _ = QFileDialog.getOpenFileName(self, "Select File to Add", "", "All Files (*);;PDF Files (*.pdf);;Word Files (*.docx);;PowerPoint Files (*.pptx);;Text Files (*.txt);;Image Files (*.jpg; *.jpeg; *.png)", options=options)
        if file_path:
            if not self.input_folder:
                QMessageBox.warning(self, "Input Folder Not Selected", "Please select an input folder first.")
                return
            if not self.output_folder:
                QMessageBox.warning(self, "Output Folder Not Selected", "Please select an output folder first.")
                return
            self.process_individual_file(file_path)

    def get_files_from_folder(self, folder):
        supported_extensions = ('.pdf', '.docx', '.pptx', '.txt', '.jpg', '.jpeg', '.png')
        files = []
        for root, _, filenames in os.walk(folder):
            for filename in filenames:
                if filename.lower().endswith(supported_extensions):
                    files.append(os.path.join(root, filename))
        logging.debug(f"Files found: {files}")
        return files

    def is_folder_categorized(self, folder):
        subfolders = [f.path for f in os.scandir(folder) if f.is_dir()]
        for subfolder in subfolders:
            if any(file.endswith(('.pdf', '.docx', '.pptx', '.txt', '.jpg', '.jpeg', '.png')) for file in os.listdir(subfolder)):
                return True
        return False

    def organize_input_folder(self):
        if self.is_folder_categorized(self.input_folder):
            QMessageBox.warning(self, "Folder Already Categorized", "The selected input folder is already categorized.")
            return
        
        files = self.get_files_from_folder(self.input_folder)
        if not files:
            QMessageBox.warning(self, "No Supported Files Found", "The selected input folder does not contain any supported files.")
            return

        self.progress_bar.setValue(0)
        self.progress_bar.setMaximum(len(files))

        categorized_files = self.categorize_files(files, self.progress_bar)
        
        for category, files in categorized_files.items():
            category_folder = os.path.join(self.output_folder, category)
            os.makedirs(category_folder, exist_ok=True)
            for file in files:
                shutil.move(file, os.path.join(category_folder, os.path.basename(file)))
        
        QMessageBox.information(self, "Organization Complete", "Files have been successfully organized.")
        self.progress_bar.setValue(len(files))

    def process_individual_file(self, file_path):
        files = [file_path]
        self.progress_bar.setValue(0)
        self.progress_bar.setMaximum(len(files))

        categorized_files = self.categorize_files(files, self.progress_bar)
        
        for category, files in categorized_files.items():
            category_folder = os.path.join(self.output_folder, category)
            os.makedirs(category_folder, exist_ok=True)
            for file in files:
                shutil.move(file, os.path.join(category_folder, os.path.basename(file)))
        
        QMessageBox.information(self, "File Processed", f"{os.path.basename(file_path)} has been successfully categorized and moved.")
        self.progress_bar.setValue(len(files))

    def extract_text_from_pdf(self, file_path):
        try:
            doc = fitz.open(file_path)
            text = ''
            for page in doc:
                text += page.get_text()
            return text
        except Exception as e:
            logging.error(f"Error extracting text from PDF: {e}")
            return ""

    def extract_text_from_docx(self, file_path):
        try:
            doc = Document(file_path)
            text = ''
            for paragraph in doc.paragraphs:
                text += paragraph.text
            return text
        except Exception as e:
            logging.error(f"Error extracting text from DOCX: {e}")
            return ""

    def extract_text_from_pptx(self, file_path):
        try:
            prs = Presentation(file_path)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            return text
        except Exception as e:
            logging.error(f"Error extracting text from PPTX: {e}")
            return ""

    def extract_text_from_txt(self, file_path):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logging.error(f"Error extracting text from TXT: {e}")
            return ""

    def extract_text_from_image(self, file_path):
        try:
            # Extract text using EasyOCR
            ocr_result = self.reader.readtext(file_path)
            ocr_text = ' '.join([res[1] for res in ocr_result])

            # Detect objects using YOLOv8
            yolo_results = self.yolo_model(file_path)
            yolo_text = ' '.join([det['name'] for det in yolo_results[0].boxes.data])

            # Combine OCR text and YOLO detected objects
            combined_text = ocr_text + ' ' + yolo_text

            logging.debug(f"Extracted text from image {file_path}: {combined_text}")
            return combined_text
        except Exception as e:
            logging.error(f"Error extracting text from image {file_path}: {e}")
            return ""

    def extract_text_from_file(self, file_path):
        _, file_extension = os.path.splitext(file_path)
        try:
            if file_extension.lower() == '.pdf':
                text = self.extract_text_from_pdf(file_path)
            elif file_extension.lower() == '.docx':
                text = self.extract_text_from_docx(file_path)
            elif file_extension.lower() == '.pptx':
                text = self.extract_text_from_pptx(file_path)
            elif file_extension.lower() == '.txt':
                text = self.extract_text_from_txt(file_path)
            elif file_extension.lower() in ['.jpg', '.jpeg', '.png']:
                text = self.extract_text_from_image(file_path)
            else:
                text = ""
            logging.debug(f"Extracted text from {file_path}: {text}")
            return text
        except Exception as e:
            logging.error(f"Error extracting text from {file_path}: {e}")
            return ""

    def categorize_files(self, files, progress_bar):
        texts = [self.extract_text_from_file(file) for file in files]
        logging.debug(f"Extracted texts: {texts}")
        
        vectorizer = TfidfVectorizer(stop_words='english')
        X = vectorizer.fit_transform(texts)

        num_clusters = min(5, len(files))  # Adjust number of clusters based on number of files
        kmeans = KMeans(n_clusters=num_clusters, random_state=42)
        labels = kmeans.fit_predict(X)

        logging.debug(f"KMeans labels: {labels}")
        
        categorized_files = defaultdict(list)
        for label, file in zip(labels, files):
            categorized_files[f"Category_{label + 1}"].append(file)

        for category, files in categorized_files.items():
            logging.debug(f"{category}: {[os.path.basename(f) for f in files]}")

        progress_bar.setValue(len(files))
        return categorized_files

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = DocumentOrganizer()
    ex.show()
    sys.exit(app.exec_())
